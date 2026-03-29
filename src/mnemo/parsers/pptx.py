"""PPTX parser — extracts slide text and metadata via python-pptx."""

from __future__ import annotations

import hashlib
from pathlib import Path

from mnemo.parser import NoteDocument

try:
    from pptx import Presentation as _Presentation  # type: ignore[import-untyped]

    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False


def parse_pptx(file_path: Path, vault_root: Path | None = None) -> NoteDocument:
    """Parse a PPTX file into a NoteDocument.

    Requires the ``python-pptx`` library (``pip install python-pptx>=0.6``).
    """
    if not _HAS_PPTX:
        raise ImportError(
            "python-pptx is required for PPTX parsing. Install with: pip install 'mnemo-secondbrain[documents]'"
        )

    prs = _Presentation(str(file_path))

    # --- extract text per slide -----------------------------------------------
    slides_text: list[str] = []
    headings: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        slide_title: str | None = None

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                slide_lines.append(text)

            # Try to identify the slide title
            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                ph_type = shape.placeholder_format.type if shape.placeholder_format else None
                # placeholder type 15 = TITLE, 13 = CENTER_TITLE
                if ph_type is not None and int(ph_type) in (15, 13, 1):
                    candidate = shape.text_frame.text.strip()
                    if candidate:
                        slide_title = candidate

        if slide_title:
            headings.append(slide_title)

        if slide_lines:
            header = f"## Slide {slide_idx}"
            if slide_title:
                header += f": {slide_title}"
            slides_text.append(header + "\n" + "\n".join(slide_lines))

    body = "\n\n".join(slides_text)

    # --- extract notes --------------------------------------------------------
    notes_parts: list[str] = []
    for slide in prs.slides:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                notes_parts.append(notes_text)
    if notes_parts:
        body = body + "\n\n## Speaker Notes\n" + "\n\n".join(notes_parts)

    # --- metadata / synthetic frontmatter ------------------------------------
    core = prs.core_properties if hasattr(prs, "core_properties") else None
    title = file_path.stem
    author = ""
    created = None
    if core:
        title = (getattr(core, "title", "") or "").strip() or file_path.stem
        author = (getattr(core, "author", "") or "").strip()
        raw_created = getattr(core, "created", None)
        if raw_created:
            created = raw_created.isoformat() if hasattr(raw_created, "isoformat") else str(raw_created)

    frontmatter: dict = {
        "title": title,
        "type": "source",
        "format": "pptx",
        "tags": ["pptx", "external"],
        "slide_count": len(prs.slides),
    }
    if author:
        frontmatter["author"] = author
    if created:
        frontmatter["created"] = created

    # --- key ------------------------------------------------------------------
    if vault_root is not None:
        try:
            rel = file_path.resolve().relative_to(vault_root.resolve())
            key = rel.with_suffix("").as_posix()
        except ValueError:
            key = file_path.stem
    else:
        key = file_path.stem

    content = body
    checksum = hashlib.sha256(content.encode("utf-8")).hexdigest()[:16]

    return NoteDocument(
        path=file_path,
        name=file_path.stem,
        key=key,
        content=content,
        body=body,
        frontmatter=frontmatter,
        wiki_links=[],
        tags=frontmatter["tags"],
        headings=headings[:30],
        checksum=checksum,
    )
